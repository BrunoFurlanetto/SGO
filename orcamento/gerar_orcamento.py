import os
import re
from copy import deepcopy
from io import BytesIO

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from orcamento.models import Tratativas, StatusOrcamento
from peraltas.models import Vendedor


class OrcamentoPDF:
    def __init__(self, id_tratativa):
        self.tratativa = Tratativas.objects.get(id_tratativa=id_tratativa)
        self.cliente = self.tratativa.cliente
        self.orcamentos = self.orcamentos_abertos
        self.responsavel = self.tratativa.orcamentos.all()[0].responsavel
        self.colaborador = self.tratativa.orcamentos.all()[0].colaborador
        self.vendedora = Vendedor.objects.get(usuario=self.colaborador)
        self.modelo_pptx = "orcamento/modelos/modelo.pptx"
        self.presentation = Presentation(self.modelo_pptx)

    def gerar_pdf(self):
        self.__alterar_dados_cliente_responsavel()
        self.__alterar_dados_vendedora()
        self.__duplicar_slide(5)
        self.trocar_posicao_slides()
        self.__substistituir_dados_pacote()
        self.__preencher_final()
        self.__salvar_pptx()
        self.__converter_pdf()

    @property
    def nome_cliente(self):
        return self.cliente.nome_fantasia

    @property
    def xml_slides(self):
        return self.presentation.slides._sldIdLst

    @property
    def n_orcamentos(self):
        if self.tratativa.orcamento_aceito:
            return 1

        status_aberto = StatusOrcamento.objects.get(
            analise_gerencia=False,
            aprovacao_cliente=False,
            negado_cliente=False,
            orcamento_vencido=False
        )

        return len(self.tratativa.orcamentos.all().filter(status_orcamento=status_aberto))

    @property
    def orcamentos_abertos(self):
        if self.tratativa.orcamento_aceito:
            return self.tratativa.orcamento_aceito

        status_aberto = StatusOrcamento.objects.get(
            analise_gerencia=False,
            aprovacao_cliente=False,
            negado_cliente=False,
            orcamento_vencido=False
        )

        return self.tratativa.orcamentos.all().filter(status_orcamento=status_aberto)

    def trocar_posicao_slides(self):
        if self.n_orcamentos > 1:
            for i, _ in enumerate(range(self.n_orcamentos - 1), start=0):
                self.__move_slide(-1, 5 + i)

    def __move_slide(self, old_index, new_index):
        slides = list(self.xml_slides)
        self.xml_slides.remove(slides[old_index])
        self.xml_slides.insert(new_index, slides[old_index])

    def __delete_slide(self, index):
        slides = list(self.xml_slides)
        self.xml_slides.remove(slides[index])

    def __duplicar_slide(self, numero_slide):
        if self.n_orcamentos <= 1:
            return

        def copiar_forma(forma_original, destino):
            copia = deepcopy(forma_original)
            destino.shapes._spTree.insert_element_before(copia._element, 'p:extLst')

        for i in range(self.n_orcamentos - 1):
            slide_original = self.presentation.slides[numero_slide - 1]
            imgDict = {}
            novo_slide = self.presentation.slides.add_slide(slide_original.slide_layout)

            for shape_original in slide_original.shapes:
                if shape_original.shape_type == 13:
                    if 'Imagem' in shape_original.name:

                        with open(shape_original.name + '.jpg', 'wb') as f:
                            f.write(shape_original.image.blob)

                        imgDict[shape_original.name + '.jpg'] = [shape_original.left, shape_original.top,
                                                                 shape_original.width, shape_original.height]
                        for k, v in imgDict.items():
                            a = novo_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
                            os.remove(k)
                else:
                    copiar_forma(shape_original, novo_slide)

        return self.presentation

    def __obter_shape_por_nome(self, nome_do_shape, slide):
        slide = self.presentation.slides[slide - 1]

        for shape in slide.shapes:
            if shape.name == nome_do_shape:
                return shape

        return None

    def __substituir_texto(
            self,
            nome_do_shape,
            valor_shape,
            fonte,
            cor_fonte,
            tamanho_fonte,
            slide,
            alinhamento=PP_PARAGRAPH_ALIGNMENT.JUSTIFY,
            negrito=False,
            italico=False,
            tachado=False,
            underline=False
    ):
        shape = self.__obter_shape_por_nome(nome_do_shape, slide)

        if shape:
            shape.text = valor_shape

            for p in shape.text_frame.paragraphs:
                for _p in p.runs:
                    _p.font.name = fonte
                    _p.font.size = tamanho_fonte
                    _p.font.color.rgb = cor_fonte
                    _p.font.bold = negrito
                    _p.font.italic = italico
                    _p.font.strike = tachado
                    _p.font.underline = underline
                    shape.text_frame.word_wrap = True
                    shape.text_frame.auto_size = True

                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = alinhamento

    def __salvar_pptx(self):
        nome_arquivo = re.sub(r"[^a-zA-Z0-9\s]", "", self.cliente.nome_fantasia)
        self.presentation.save('temp\\orcamento.pptx')

    @staticmethod
    def __converter_pdf():
        from spire.presentation import Presentation as Prs, FileFormat

        # Create a Presentation object
        presentation = Prs()
        # Load a PowerPoint presentation in PPTX format
        presentation.LoadFromFile("temp/orcamento.pptx")
        # Or load a PowerPoint presentation in PPT format
        # presentation.LoadFromFile("Sample.ppt")

        # Convert the presentation to PDF format
        presentation.SaveToFile("temp/orcamento.pdf", FileFormat.PDF)
        presentation.Dispose()

    def __alterar_dados_cliente_responsavel(self):
        fonte = 'IBM Plex Sans'
        cor_fonte = RGBColor(131, 128, 131)
        tamanho_fonte = Pt(14)
        self.__substituir_texto('NomeColegio', self.cliente.nome_fantasia, fonte, cor_fonte, tamanho_fonte, 4)
        self.__substituir_texto('NomeResponsavel', self.responsavel.nome, fonte, cor_fonte, tamanho_fonte, 4)
        self.__substituir_texto('TelefoneResponsavel', self.responsavel.fone, fonte, cor_fonte, tamanho_fonte, 4)
        self.__substituir_texto('EmailResponsavel', self.responsavel.email_responsavel_evento, fonte, cor_fonte,
                                tamanho_fonte, 4)

    def __alterar_dados_vendedora(self):
        fonte = 'IBM Plex Sans'
        cor_fonte = RGBColor(131, 128, 131)
        tamanho_fonte = Pt(14)
        self.__substituir_texto('NomeVendedora', self.colaborador.get_full_name(), fonte, cor_fonte, tamanho_fonte, 4)
        self.__substituir_texto('TelefoneVendedora', self.vendedora.telefone, fonte, cor_fonte, tamanho_fonte, 4)
        self.__substituir_texto('WhatsVendedora', self.vendedora.telefone, fonte, cor_fonte, tamanho_fonte, 4)
        self.__substituir_texto('EmailVendedora', self.colaborador.email, fonte, cor_fonte, tamanho_fonte, 4)

    def __substistituir_dados_pacote(self):
        if self.n_orcamentos > 1:
            for i, orcamento in enumerate(self.orcamentos, start=0):
                fonte = 'Calibri (Corpo)'
                cor_fonte = RGBColor(0, 0, 0)
                tamanho_fonte = Pt(14)
                texto_ceu = f'{orcamento.produto.produto} + 0{len(orcamento.atividades_ceu.all())} AT. CEU'
                texto_valor = f'Por R$ {orcamento.valor} por participante'
                texto_entrada = f'Chegada às \n{orcamento.check_in.time()} \nSaída às \n{orcamento.check_out.time()}'
                texto_incluso = f'- Café,  almoço  e lanche antes da saída \n- 01 monitor de lazer a cada 20 pagantes \n- {len(orcamento.atividades_ceu.all())} atividade na Fundação CEU (exceto oficina de foguetes) \n- Programação de lazer.'
                texto_final = f'\u2713    À vista R$ {orcamento.valor} ou, parcelado via link de pagamento online. Consulte quantidade de parcelas.'

                self.__substituir_texto('ProdutoEceu', texto_ceu, fonte, cor_fonte, tamanho_fonte, 5 + i, negrito=True,
                                        alinhamento=PP_PARAGRAPH_ALIGNMENT.CENTER)
                self.__substituir_texto('ValorParticipante', texto_valor, fonte, cor_fonte, tamanho_fonte, 5 + i,
                                        negrito=True, alinhamento=PP_PARAGRAPH_ALIGNMENT.CENTER)
                self.__substituir_texto('DadosEntradaSaida', texto_entrada, fonte, cor_fonte, tamanho_fonte, 5 + i,
                                        negrito=True, alinhamento=PP_PARAGRAPH_ALIGNMENT.CENTER)
                self.__substituir_texto('DadosIncluso', texto_incluso, fonte, cor_fonte, tamanho_fonte, 5 + i)
                self.__substituir_texto('TextoFinalValor', texto_final, fonte, cor_fonte, Pt(12), 5 + i, negrito=True)
        else:
            orcamento = self.tratativa.orcamento_aceito
            fonte = 'Calibri (Corpo)'
            cor_fonte = RGBColor(0, 0, 0)
            tamanho_fonte = Pt(14)
            texto_ceu = f'{orcamento.produto.produto} + 0{len(orcamento.atividades_ceu.all())} AT. CEU'
            texto_valor = f'Por R$ {str(orcamento.valor).replace(".", ",")} por participante'
            texto_entrada = f'Chegada às \n{orcamento.check_in.time()} \nSaída às \n{orcamento.check_out.time()}'
            texto_incluso = f'- Café,  almoço  e lanche antes da saída \n- 01 monitor de lazer a cada 20 pagantes \n- {len(orcamento.atividades_ceu.all())} atividade na Fundação CEU (exceto oficina de foguetes) \n- Programação de lazer.'
            texto_final = f'\u2713    À vista R$ {str(orcamento.valor).replace(".", ",")} ou, parcelado via link de pagamento online. Consulte quantidade de parcelas.'

            self.__substituir_texto('ProdutoEceu', texto_ceu, fonte, cor_fonte, tamanho_fonte, 5, negrito=True,
                                    alinhamento=PP_PARAGRAPH_ALIGNMENT.CENTER)
            self.__substituir_texto('ValorParticipante', texto_valor, fonte, cor_fonte, tamanho_fonte, 5,
                                    negrito=True, alinhamento=PP_PARAGRAPH_ALIGNMENT.CENTER)
            self.__substituir_texto('DadosEntradaSaida', texto_entrada, fonte, cor_fonte, tamanho_fonte, 5,
                                    negrito=True, alinhamento=PP_PARAGRAPH_ALIGNMENT.CENTER)
            self.__substituir_texto('DadosIncluso', texto_incluso, fonte, cor_fonte, tamanho_fonte, 5)
            self.__substituir_texto('TextoFinalValor', texto_final, fonte, cor_fonte, Pt(12), 5, negrito=True)

    def __preencher_final(self):
        fonte = 'IBM Plex Sans'
        cor_fonte = RGBColor(131, 128, 131)
        n_slides = len(self.presentation.slides)

        self.__substituir_texto('NomeVendedora', self.colaborador.get_full_name(), fonte, cor_fonte, Pt(18), n_slides)
        self.__substituir_texto('TelefoneVendedora', self.vendedora.telefone, fonte, cor_fonte, Pt(11), n_slides)
        self.__substituir_texto('WhatsVendedora', self.vendedora.telefone, fonte, cor_fonte, Pt(11), n_slides)
        self.__substituir_texto('EmailVendedora', self.colaborador.email, fonte, cor_fonte, Pt(12), n_slides)
